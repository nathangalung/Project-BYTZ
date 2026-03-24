"""CV text parser with string matching algorithms."""

import re
from dataclasses import dataclass, field

# Skill taxonomy
SKILL_DB: list[str] = [
    "React", "Vue.js", "Angular", "Next.js", "TypeScript", "JavaScript",
    "HTML", "CSS", "Tailwind CSS", "Bootstrap", "SASS",
    "Node.js", "Express", "Hono", "NestJS", "Fastify",
    "Python", "Django", "Flask", "FastAPI",
    "Go", "Golang", "Fiber", "Gin",
    "Java", "Spring Boot", "Kotlin",
    "PHP", "Laravel", "CodeIgniter",
    "Ruby", "Ruby on Rails",
    "Rust", "C++", "C#", ".NET",
    "React Native", "Flutter", "Swift", "Dart", "Expo",
    "PostgreSQL", "MySQL", "MongoDB", "Redis", "SQLite",
    "Docker", "Kubernetes", "AWS", "GCP", "Azure",
    "CI/CD", "GitHub Actions", "Terraform", "Linux",
    "Figma", "Adobe XD", "Sketch", "Photoshop", "Illustrator",
    "TensorFlow", "PyTorch", "Pandas", "Scikit-learn",
    "Git", "REST API", "GraphQL", "gRPC",
    "Agile", "Scrum", "Jira",
]

# Skill aliases for fuzzy matching
SKILL_ALIASES: dict[str, str] = {
    "reactjs": "React", "react.js": "React",
    "vuejs": "Vue.js", "vue": "Vue.js",
    "angularjs": "Angular",
    "ts": "TypeScript", "typescript": "TypeScript",
    "js": "JavaScript", "javascript": "JavaScript",
    "nodejs": "Node.js", "node": "Node.js",
    "expressjs": "Express", "express.js": "Express",
    "golang": "Go",
    "postgres": "PostgreSQL", "psql": "PostgreSQL",
    "mongo": "MongoDB",
    "k8s": "Kubernetes",
    "tailwind": "Tailwind CSS",
    "nextjs": "Next.js", "next": "Next.js",
    "nestjs": "NestJS",
    "springboot": "Spring Boot",
    "ruby on rails": "Ruby on Rails", "rails": "Ruby on Rails",
    "react native": "React Native",
    "github actions": "GitHub Actions",
    "rest": "REST API", "restful": "REST API",
    "graphql": "GraphQL",
    "tensorflow": "TensorFlow",
    "pytorch": "PyTorch",
    "scikit-learn": "Scikit-learn", "sklearn": "Scikit-learn",
    "ci/cd": "CI/CD", "cicd": "CI/CD",
}


def levenshtein_distance(s1: str, s2: str) -> int:
    """Levenshtein distance for fuzzy matching."""
    if len(s1) < len(s2):
        return levenshtein_distance(s2, s1)
    if len(s2) == 0:
        return len(s1)
    prev = list(range(len(s2) + 1))
    for i, c1 in enumerate(s1):
        curr = [i + 1]
        for j, c2 in enumerate(s2):
            cost = 0 if c1 == c2 else 1
            curr.append(min(curr[j] + 1, prev[j + 1] + 1, prev[j] + cost))
        prev = curr
    return prev[-1]


class AhoCorasick:
    """Multi-pattern matching."""

    def __init__(self) -> None:
        self.goto: list[dict[str, int]] = [{}]
        self.fail: list[int] = [0]
        self.output: list[list[str]] = [[]]

    def add_pattern(self, pattern: str, label: str) -> None:
        state = 0
        for ch in pattern.lower():
            if ch not in self.goto[state]:
                self.goto[state][ch] = len(self.goto)
                self.goto.append({})
                self.fail.append(0)
                self.output.append([])
            state = self.goto[state][ch]
        self.output[state].append(label)

    def build(self) -> None:
        from collections import deque
        q: deque[int] = deque()
        for ch, s in self.goto[0].items():
            q.append(s)
        while q:
            r = q.popleft()
            for ch, s in self.goto[r].items():
                q.append(s)
                state = self.fail[r]
                while state != 0 and ch not in self.goto[state]:
                    state = self.fail[state]
                self.fail[s] = self.goto[state].get(ch, 0)
                if self.fail[s] == s:
                    self.fail[s] = 0
                self.output[s] = self.output[s] + self.output[self.fail[s]]

    def search(self, text: str) -> set[str]:
        results: set[str] = set()
        state = 0
        for ch in text.lower():
            while state != 0 and ch not in self.goto[state]:
                state = self.fail[state]
            state = self.goto[state].get(ch, 0)
            for label in self.output[state]:
                results.add(label)
        return results


def build_skill_matcher() -> AhoCorasick:
    """Build Aho-Corasick automaton for skills."""
    ac = AhoCorasick()
    for skill in SKILL_DB:
        ac.add_pattern(skill.lower(), skill)
    for alias, canonical in SKILL_ALIASES.items():
        ac.add_pattern(alias, canonical)
    ac.build()
    return ac


_SKILL_MATCHER = build_skill_matcher()


def extract_skills_from_text(text: str) -> list[str]:
    """Extract skills using Aho-Corasick + Levenshtein fallback."""
    # Phase 1: Aho-Corasick exact + alias matching
    found = _SKILL_MATCHER.search(text)

    # Phase 2: Levenshtein fuzzy match on remaining words
    words = re.findall(r'\b[A-Za-z][A-Za-z.#+/\-]{1,20}\b', text)
    for word in words:
        w_lower = word.lower()
        if any(w_lower == s.lower() for s in found):
            continue
        for skill in SKILL_DB:
            if levenshtein_distance(w_lower, skill.lower()) <= 2 and len(skill) > 3:
                found.add(skill)
                break

    return sorted(found)


def extract_text(file_bytes: bytes, file_type: str) -> str:
    """Extract raw text from a document given its bytes and type."""
    import tempfile
    from pathlib import Path

    ext = file_type.lower()
    text = ""

    try:
        if ext == "pdf":
            try:
                import pypdfium2 as pdfium

                pdf = pdfium.PdfDocument(file_bytes)
                pages = []
                for page in pdf:
                    textpage = page.get_textpage()
                    pages.append(textpage.get_text_bounded())
                    textpage.close()
                    page.close()
                pdf.close()
                text = "\n".join(pages)
            except Exception:
                text = file_bytes.decode("utf-8", errors="ignore")
        elif ext in ("docx", "doc"):
            try:
                import docx

                with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
                    tmp.write(file_bytes)
                    tmp_path = tmp.name
                doc = docx.Document(tmp_path)
                text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
                Path(tmp_path).unlink(missing_ok=True)
            except Exception:
                text = file_bytes.decode("utf-8", errors="ignore")
        elif ext == "pptx":
            try:
                from pptx import Presentation

                with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
                    tmp.write(file_bytes)
                    tmp_path = tmp.name
                prs = Presentation(tmp_path)
                parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            parts.append(shape.text)
                text = "\n".join(parts)
                Path(tmp_path).unlink(missing_ok=True)
            except Exception:
                text = file_bytes.decode("utf-8", errors="ignore")
        else:
            text = file_bytes.decode("utf-8", errors="ignore")
    except Exception:
        text = file_bytes.decode("utf-8", errors="ignore")

    return text


@dataclass
class ParsedCV:
    """Structured CV data."""
    name: str = ""
    email: str = ""
    phone: str = ""
    skills: list[str] = field(default_factory=list)
    education: list[dict[str, str]] = field(default_factory=list)
    experience: list[dict[str, str]] = field(default_factory=list)
    projects: list[dict[str, str]] = field(default_factory=list)
    portfolio_urls: list[str] = field(default_factory=list)


def extract_emails(text: str) -> list[str]:
    return re.findall(r'[\w.+-]+@[\w-]+\.[\w.-]+', text)


def extract_phones(text: str) -> list[str]:
    return re.findall(r'(?:\+62|62|0)\d[\d\s\-]{8,14}', text)


def extract_urls(text: str) -> list[str]:
    urls = re.findall(r'https?://[^\s<>"\']+', text)
    portfolio_domains = ['github.com', 'linkedin.com', 'dribbble.com', 'behance.net', 'gitlab.com']
    return [u for u in urls if any(d in u for d in portfolio_domains)]


def extract_name_heuristic(text: str) -> str:
    """Extract name from first non-empty line."""
    for line in text.strip().split('\n'):
        line = line.strip()
        if 2 < len(line) < 50 and not re.search(r'[@\d]', line) and not line.startswith('http'):
            return line
    return ""


def parse_cv_text(text: str) -> ParsedCV:
    """Parse raw CV text into structured data."""
    result = ParsedCV()
    result.name = extract_name_heuristic(text)
    emails = extract_emails(text)
    if emails:
        result.email = emails[0]
    phones = extract_phones(text)
    if phones:
        result.phone = phones[0].strip()
    result.skills = extract_skills_from_text(text)
    result.portfolio_urls = extract_urls(text)

    # Extract education (simple heuristic)
    edu_patterns = [
        r'(?:universitas|university|institut|politeknik)\s+[\w\s]+',
        r'(?:S1|S2|S3|Bachelor|Master|PhD)\s+[\w\s]+',
    ]
    for pat in edu_patterns:
        matches = re.findall(pat, text, re.IGNORECASE)
        for m in matches[:2]:
            result.education.append({"university": m.strip()})

    # Extract experience (simple heuristic)
    exp_patterns = re.findall(
        r'(\d{4})\s*[-–]\s*(\d{4}|present|sekarang)[\s:]*(.+)',
        text, re.IGNORECASE
    )
    for start, end, desc in exp_patterns[:5]:
        result.experience.append({
            "position": desc.strip()[:100],
            "start": start,
            "end": end,
        })

    return result
